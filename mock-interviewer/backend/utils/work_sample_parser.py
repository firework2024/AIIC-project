from pathlib import Path

from PyPDF2 import PdfReader


SUPPORTED_WORK_SAMPLE_EXTENSIONS = {".pdf", ".txt", ".pptx"}


def extract_text_from_pdf(file) -> str:
    reader = PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
        text += "\n"
    return text


def extract_text_from_txt(file) -> str:
    raw = file.read()
    if isinstance(raw, str):
        return raw
    for encoding in ("utf-8", "gbk", "utf-16"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="ignore")


def extract_text_from_pptx(file) -> str:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise RuntimeError("服务器暂未安装 PPTX 解析依赖。") from exc

    presentation = Presentation(file)
    chunks = []
    for page_number, slide in enumerate(presentation.slides, start=1):
        slide_chunks = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                slide_chunks.append(text.strip())
        if slide_chunks:
            chunks.append(f"第 {page_number} 页：\n" + "\n".join(slide_chunks))
    return "\n\n".join(chunks)


def extract_work_sample(file):
    filename = file.filename or ""
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_WORK_SAMPLE_EXTENSIONS:
        return False, "请上传 PDF、TXT 或 PPTX 格式的作品。"

    try:
        if suffix == ".pdf":
            text = extract_text_from_pdf(file)
        elif suffix == ".txt":
            text = extract_text_from_txt(file)
        else:
            text = extract_text_from_pptx(file)
    except Exception:
        return False, "作品文件解析失败，请换成文字可复制的 PDF、TXT 或 PPTX。"

    text = (text or "").strip()
    if len(text) < 200:
        return False, "作品内容太短或无法提取有效文字，请上传内容更完整的报告、PPT 或文本。"

    return True, text[:12000]
